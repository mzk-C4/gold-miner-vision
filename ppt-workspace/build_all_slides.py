"""
Generate complete 10-slide PPTX presentation for Gold Miner.
All images embedded as base64, no external dependencies.
Each slide is rendered via Playwright HTML screenshot → embedded in PPTX.
"""
import base64
import io
import os
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from PIL import Image

WORKSPACE = Path(r"D:\GOLD MINER\ppt-workspace")
ASSETS_DIR = WORKSPACE / "html" / "assets"
OUTPUT_PPTX = WORKSPACE / "pptx" / "黄金矿工-手势识别版-完整版.pptx"

SLIDE_W = 1920
SLIDE_H = 1080
PPTX_W = Inches(13.333)
PPTX_H = Inches(7.5)

# ============================================================
# Load game sprites as base64
# ============================================================
def img_b64(path):
    """Load image and return base64 data URI."""
    with open(path, 'rb') as f:
        data = base64.b64encode(f.read()).decode()
    # Detect mime type
    ext = Path(path).suffix.lower()
    mime = 'image/png' if ext == '.png' else 'image/jpeg'
    return f'data:{mime};base64,{data}'

# Load all game sprites
SPRITES = {}
for png_file in sorted(ASSETS_DIR.glob('*.png')):
    SPRITES[png_file.stem] = img_b64(png_file)
    print(f"  Loaded sprite: {png_file.stem}")

# ============================================================
# CSS Framework (shared across all slides)
# ============================================================
CSS_COMMON = """
* { margin:0; padding:0; box-sizing:border-box; }
body {
  width:1920px; height:1080px; overflow:hidden;
  background:#0A0E27; font-family:'Microsoft YaHei','PingFang SC','Noto Sans SC',sans-serif;
  display:flex; align-items:center; justify-content:center;
}
.slide { width:1920px; height:1080px; position:relative; background:#0A0E27; overflow:hidden; }
.slide::before {
  content:''; position:absolute; inset:0;
  background: radial-gradient(ellipse at center, transparent 40%, rgba(10,14,39,0.5) 100%);
  pointer-events:none; z-index:1;
}
"""

# ============================================================
# Helper: render HTML → PNG bytes
# ============================================================
def render_slide(html_content, browser):
    """Render a full HTML page to PNG screenshot bytes."""
    page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
    page.set_content(html_content)
    page.wait_for_timeout(1000)
    screenshot = page.screenshot(full_page=False)
    page.close()
    return screenshot

# ============================================================
# SLIDE 01: Cover
# ============================================================
def slide_01():
    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.main-title {{
  font-family: 'Georgia','SimSun',serif;
  font-size:130px; font-weight:900; letter-spacing:14px;
  background: linear-gradient(180deg,#F0D68A 0%,#D4A843 30%,#B8860B 55%,#E8C870 70%,#D4A843 100%);
  -webkit-background-clip:text; -webkit-text-fill-color:transparent;
  background-clip:text;
  filter: drop-shadow(0 4px 10px rgba(0,0,0,0.6));
  line-height:1.1; margin-bottom:6px;
}}
.sub-title {{ font-size:44px; font-weight:300; color:#E0DDD8; letter-spacing:18px; margin-bottom:30px; }}
.sep {{ width:200px; height:1px; background:linear-gradient(90deg,transparent,#A08050,#D4A843,#A08050,transparent); margin:0 auto 30px; opacity:0.7; }}
.course {{ font-size:22px; font-weight:300; color:#8A8FA8; letter-spacing:10px; }}
.title-area {{ position:absolute; top:290px; left:50%; transform:translateX(-50%); text-align:center; z-index:6; }}
.hook-rope {{ position:absolute; top:-10px; left:50%; transform:translateX(-50%); z-index:5; }}
.rope {{ width:3px; height:190px; margin:0 auto; background:linear-gradient(to bottom,#A08060,#C4A87C,#8B7355,#B8946E,#9B7D5C,#C4A87C); }}
.hand-box {{ position:absolute; right:70px; top:50%; transform:translateY(-50%); z-index:2; opacity:0.28; }}
.nuggets {{ position:absolute; bottom:50px; left:50%; transform:translateX(-50%); display:flex; gap:14px; align-items:flex-end; z-index:6; }}
.nuggets img {{ display:block; filter:drop-shadow(0 4px 14px rgba(0,0,0,0.6)); }}
.spot {{ position:absolute; bottom:0; left:50%; transform:translateX(-50%); width:750px; height:180px;
  background:radial-gradient(ellipse at center bottom, rgba(212,168,67,0.07) 0%, transparent 70%); z-index:3; }}
.top-glow {{ position:absolute; top:-80px; left:50%; transform:translateX(-50%); width:800px; height:350px;
  background:radial-gradient(ellipse at center top, rgba(200,180,140,0.04) 0%, transparent 70%); z-index:3; }}
</style></head><body>
<div class="slide">
<div class="top-glow"></div>
<div class="hook-rope">
  <div class="rope"></div>
  <svg width="130" height="150" viewBox="0 0 130 150">
    <defs><linearGradient id="hg" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0%" stop-color="#D4A843"/><stop offset="35%" stop-color="#F0D88A"/>
      <stop offset="60%" stop-color="#B8860B"/><stop offset="100%" stop-color="#8B6914"/>
    </linearGradient></defs>
    <circle cx="65" cy="18" r="9" fill="none" stroke="url(#hg)" stroke-width="3"/>
    <path d="M65 28 L65 60 Q65 90 45 105 Q32 115 38 132 Q44 142 60 136 Q70 132 67 118 Q63 104 48 96" fill="none" stroke="url(#hg)" stroke-width="6" stroke-linecap="round"/>
    <path d="M38 132 L30 146" fill="none" stroke="url(#hg)" stroke-width="5" stroke-linecap="round"/>
    <path d="M63 35 L63 58 Q63 82 47 97 Q38 108 42 122" fill="none" stroke="rgba(255,255,255,0.12)" stroke-width="1.5"/>
  </svg>
</div>
<div class="title-area">
  <div class="main-title">黄金矿工</div>
  <div class="sub-title">手势识别版</div>
  <div class="sep"></div>
  <div class="course">C++ 课程设计汇报</div>
</div>
<div class="hand-box">
  <svg width="320" height="420" viewBox="0 0 320 420">
    <defs>
      <linearGradient id="hnd" x1="0" y1="0" x2="1" y2="1">
        <stop offset="0%" stop-color="#4A7FB5" stop-opacity="0.5"/>
        <stop offset="100%" stop-color="#3A6090" stop-opacity="0.1"/>
      </linearGradient>
    </defs>
    <path d="M120 210 Q100 155 110 105 Q115 75 130 65 Q140 60 145 70 L145 145 Q148 105 152 58 Q155 33 165 28 Q175 25 178 43 L175 135 Q180 88 185 48 Q188 23 198 21 Q208 23 205 63 L198 145 Q205 105 215 68 Q220 43 228 45 Q236 53 230 98 L210 175 Q200 215 180 245 Q160 265 140 265 Q120 263 115 245 Q108 225 120 210Z" fill="url(#hnd)" opacity="0.7"/>
    <path d="M115 245 Q110 275 105 315 Q108 345 120 365 Q140 375 160 370 Q175 360 178 335 Q175 295 170 265" fill="url(#hnd)" opacity="0.2"/>
    <g opacity="0.4" stroke="#4A7FB5" stroke-width="0.8" fill="none">
      <path d="M145 70 L155 45 Q158 28 162 18" stroke-dasharray="3,3"/><circle cx="162" cy="18" r="2.5" fill="rgba(74,127,181,0.5)"/>
      <path d="M178 43 L182 23 Q185 10 190 4" stroke-dasharray="3,3"/><circle cx="190" cy="4" r="2.5" fill="rgba(74,127,181,0.5)"/>
      <path d="M205 63 L210 45 Q213 30 216 20" stroke-dasharray="3,3"/><circle cx="216" cy="20" r="2.5" fill="rgba(74,127,181,0.5)"/>
      <path d="M230 98 L235 80 Q238 64 240 52" stroke-dasharray="3,3"/><circle cx="240" cy="52" r="2.5" fill="rgba(74,127,181,0.5)"/>
      <path d="M130 65 L115 58 Q100 52 90 57" stroke-dasharray="3,3"/><circle cx="90" cy="57" r="2.5" fill="rgba(74,127,181,0.5)"/>
    </g>
  </svg>
</div>
<div class="spot"></div>
<div class="nuggets">
  <img src="{SPRITES['gold-large-2']}" style="height:100px" alt="Gold">
  <img src="{SPRITES['diamond-2']}" style="height:70px" alt="Diamond">
  <img src="{SPRITES['gold-nugget-4']}" style="height:78px" alt="Gold">
  <img src="{SPRITES['treasure']}" style="height:72px" alt="Treasure">
  <img src="{SPRITES['stone-2']}" style="height:54px" alt="Stone">
  <img src="{SPRITES['gold-nugget-3']}" style="height:60px" alt="Gold">
  <img src="{SPRITES['gen-diamond']}" style="height:68px" alt="Diamond">
  <img src="{SPRITES['gold-large-1']}" style="height:85px" alt="Gold">
  <img src="{SPRITES['tnt']}" style="height:65px" alt="TNT">
</div>
</div></body></html>"""

# ============================================================
# SLIDE 02: Background
# ============================================================
def slide_02():
    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.title {{ position:absolute; top:55px; left:50%; transform:translateX(-50%); font-size:42px; font-weight:700; color:#E8E5E0; letter-spacing:8px; z-index:6; }}
.title::after {{ content:''; display:block; width:80px; height:2px; background:#D4A843; margin:12px auto 0; opacity:0.6; }}
.panel {{ position:absolute; top:180px; width:420px; z-index:5; }}
.panel-l {{ left:120px; }}
.panel-r {{ right:120px; }}
.panel-label {{ font-size:22px; font-weight:500; margin-bottom:18px; letter-spacing:4px; }}
.pl-gold {{ color:#D4A843; }} .pl-blue {{ color:#4A7FB5; }}
.panel-box {{ width:420px; height:540px; border-radius:8px; padding:30px; }}
.box-left {{ background:linear-gradient(135deg, rgba(212,168,67,0.08), rgba(180,140,60,0.03)); border:1px solid rgba(212,168,67,0.15); }}
.box-right {{ background:linear-gradient(135deg, rgba(74,127,181,0.08), rgba(50,90,130,0.03)); border:1px solid rgba(74,127,181,0.15); }}
.mine-scene {{ width:100%; height:340px; display:flex; align-items:center; justify-content:center; margin-bottom:20px; }}
.tech-scene {{ width:100%; height:340px; display:flex; align-items:center; justify-content:center; margin-bottom:20px; }}
.sep-line {{ position:absolute; left:50%; top:200px; width:1px; height:500px; background:linear-gradient(to bottom,#D4A843,#4A7FB5); transform:translateX(-50%); z-index:4; opacity:0.4; }}
.sep-x {{ position:absolute; left:50%; top:430px; transform:translate(-50%,-50%); font-size:36px; color:#D4A843; z-index:5; font-weight:300; }}
.tags {{ position:absolute; bottom:90px; left:50%; transform:translateX(-50%); display:flex; gap:20px; z-index:6; }}
.tag {{ padding:8px 22px; border-radius:4px; font-size:17px; letter-spacing:3px; background:transparent; }}
.tag-cpp {{ border:1px solid #5BAF8A; color:#5BAF8A; }}
.tag-cocos {{ border:1px solid #4A7FB5; color:#4A7FB5; }}
.tag-ocv {{ border:1px solid #D4A843; color:#D4A843; }}
.underground {{ text-align:center; }}
.underground svg {{ display:block; margin:0 auto; }}
.panel-desc {{ text-align:center; font-size:17px; color:#8A8FA8; line-height:1.8; }}
</style></head><body>
<div class="slide">
<div class="title">项目背景</div>
<div class="sep-line"></div>
<div class="sep-x">×</div>
<div class="panel panel-l">
  <div class="panel-label pl-gold">经典游戏</div>
  <div class="panel-box box-left">
    <div class="mine-scene">
      <svg width="300" height="300" viewBox="0 0 300 300">
        <rect x="50" y="20" width="200" height="180" fill="rgba(139,90,43,0.3)" rx="8"/>
        <rect x="60" y="200" width="180" height="80" fill="rgba(100,60,30,0.4)" rx="4"/>
        <circle cx="120" cy="100" r="18" fill="rgba(212,168,67,0.7)"/>
        <circle cx="180" cy="140" r="14" fill="rgba(212,168,67,0.5)"/>
        <circle cx="100" cy="160" r="10" fill="rgba(180,140,80,0.4)"/>
        <polygon points="200,90 210,75 215,90 210,80" fill="rgba(100,200,200,0.6)"/>
        <polygon points="140,50 148,35 152,50 148,42" fill="rgba(200,100,100,0.5)"/>
        <line x1="150" y1="20" x2="150" y2="70" stroke="rgba(180,150,100,0.5)" stroke-width="2"/>
        <circle cx="150" cy="75" r="6" fill="none" stroke="rgba(180,150,100,0.6)" stroke-width="2"/>
      </svg>
    </div>
    <div class="panel-desc">经典黄金矿工游戏机制<br>钩爪抓取 · 矿物收集 · 关卡挑战</div>
  </div>
</div>
<div class="panel panel-r">
  <div class="panel-label pl-blue">手势识别</div>
  <div class="panel-box box-right">
    <div class="tech-scene">
      <svg width="280" height="300" viewBox="0 0 280 300">
        <path d="M100 180 Q85 135 95 90 Q100 62 112 54 Q120 50 124 58 L124 120 Q126 88 130 48 Q133 28 140 24 Q148 22 150 36 L148 112 Q152 74 156 40 Q158 20 166 18 Q174 20 172 52 L168 120 Q172 88 180 56 Q184 36 190 38 Q196 44 192 80 L178 142 Q168 178 152 200 Q136 218 120 218 Q104 216 100 200 Q94 186 100 170Z" fill="rgba(74,127,181,0.3)"/>
        <circle cx="124" cy="58" r="4" fill="rgba(74,127,181,0.7)"/>
        <circle cx="150" cy="36" r="4" fill="rgba(74,127,181,0.7)"/>
        <circle cx="172" cy="52" r="4" fill="rgba(74,127,181,0.7)"/>
        <circle cx="192" cy="80" r="4" fill="rgba(74,127,181,0.7)"/>
        <circle cx="100" cy="100" r="4" fill="rgba(74,127,181,0.7)"/>
        <line x1="124" y1="58" x2="140" y2="100" stroke="rgba(74,127,181,0.2)" stroke-width="0.5"/>
        <line x1="150" y1="36" x2="155" y2="90" stroke="rgba(74,127,181,0.2)" stroke-width="0.5"/>
        <line x1="172" y1="52" x2="170" y2="100" stroke="rgba(74,127,181,0.2)" stroke-width="0.5"/>
        <line x1="192" y1="80" x2="180" y2="120" stroke="rgba(74,127,181,0.2)" stroke-width="0.5"/>
      </svg>
    </div>
    <div class="panel-desc">计算机视觉手势识别<br>HSV 肤色分割 · 凸包分析 · 指尖定位</div>
  </div>
</div>
<div class="tags">
  <div class="tag tag-cpp">C++17</div>
  <div class="tag tag-cocos">Cocos2d-x</div>
  <div class="tag tag-ocv">OpenCV</div>
</div>
</div></body></html>"""

# ============================================================
# SLIDE 03: Core Data Structures (KEY)
# ============================================================
def slide_03():
    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.title {{ position:absolute; top:50px; left:50%; transform:translateX(-50%); font-size:42px; font-weight:700; color:#E8E5E0; letter-spacing:6px; z-index:6; }}
.title::after {{ content:''; display:block; width:60px; height:2px; background:#D4A843; margin:10px auto 0; }}
.cards {{ position:absolute; top:160px; left:50%; transform:translateX(-50%); display:flex; gap:36px; z-index:5; }}
.card {{ width:310px; background:rgba(30,35,64,0.92); border-radius:6px; padding:24px 22px; box-shadow:0 6px 24px rgba(0,0,0,0.4); }}
.card-tab {{ height:3px; width:60%; margin:0 auto 16px; border-radius:2px; }}
.tab-gold {{ background:#D4A843; }} .tab-mint {{ background:#5BAF8A; }} .tab-blue {{ background:#4A7FB5; }}
.card-title {{ font-size:26px; font-weight:700; text-align:center; margin-bottom:4px; }}
.ct-gold {{ color:#D4A843; }} .ct-mint {{ color:#5BAF8A; }} .ct-blue {{ color:#4A7FB5; }}
.card-sub {{ font-size:14px; color:#6A6F88; text-align:center; margin-bottom:14px; }}
.card-code {{ font-family:'Consolas','Courier New',monospace; font-size:16px; line-height:2.2; color:#B0B8C8; }}
.code-key {{ color:#4A7FB5; }} .code-type {{ color:#5BAF8A; }} .code-cn {{ color:#C0C4D0; }}
.arrow {{ position:absolute; top:300px; z-index:4; font-size:28px; color:#8B7A4A; }}
.arrow-l {{ left:420px; }} .arrow-r {{ right:420px; }}
.footnote {{ position:absolute; bottom:60px; left:50%; transform:translateX(-50%); font-size:14px; color:#5A6080; z-index:5; }}
</style></head><body>
<div class="slide">
<div class="title">核心数据结构 ⭐</div>
<div class="cards" style="position:absolute; top:150px; left:50%; transform:translateX(-50%); display:flex; gap:36px;">
  <div class="card">
    <div class="card-tab tab-gold"></div>
    <div class="card-title ct-gold">GestureType</div>
    <div class="card-sub">枚举类型</div>
    <div class="card-code">
      <span class="code-key">OPEN_PALM</span> <span class="code-cn">→ 瞄准</span><br>
      <span class="code-key">FIST</span> <span class="code-cn">→ 释放</span><br>
      <span class="code-key">UNKNOWN</span> <span class="code-cn">→ 无效</span>
    </div>
  </div>
  <div style="font-size:32px; color:#8B7A4A; padding-top:120px;">→</div>
  <div class="card">
    <div class="card-tab tab-mint"></div>
    <div class="card-title ct-mint">GestureResult</div>
    <div class="card-sub">结构体</div>
    <div class="card-code">
      <span class="code-key">gesture</span><span class="code-cn"> : </span><span class="code-type">GestureType</span><br>
      <span class="code-key">angle</span><span class="code-cn"> : </span><span class="code-type">float</span><br>
      <span class="code-key">confidence</span><span class="code-cn"> : </span><span class="code-type">float</span><br>
      <span class="code-key">isValid</span><span class="code-cn"> : </span><span class="code-type">bool</span>
    </div>
  </div>
  <div style="font-size:32px; color:#8B7A4A; padding-top:120px;">→</div>
  <div class="card">
    <div class="card-tab tab-blue"></div>
    <div class="card-title ct-blue">GestureCommand</div>
    <div class="card-sub">输出指令</div>
    <div class="card-code">
      <span class="code-key">shouldReleaseHook</span><br>
      <span class="code-key">shouldDetonateBomb</span><br>
      <span class="code-key">targetAngle</span><span class="code-cn"> : </span><span class="code-type">float</span>
    </div>
  </div>
</div>
<div class="footnote">GestureType → GestureResult → GestureCommand &nbsp;|&nbsp; 数据流自左向右</div>
</div></body></html>"""

# ============================================================
# SLIDE 04: CV Pipeline (KEY)
# ============================================================
def slide_04():
    nodes = [
        ("#4A7FB5", "HSV 分割", "肤色提取"),
        ("#4A7FB5", "轮廓检测", "最大连通域"),
        ("#4A7FB5", "凸包分析", "缺陷点定位"),
        ("#5BAF8A", "指尖识别", "凸缺陷距离判定"),
        ("#D4A843", "角度计算", "手势分类输出"),
    ]
    node_html = ""
    for i, (color, label, desc) in enumerate(nodes):
        left = 110 + i * 340
        node_html += f"""
<div style="position:absolute; top:300px; left:{left}px; text-align:center; z-index:5;">
  <svg width="100" height="100" viewBox="0 0 100 100">
    <circle cx="50" cy="50" r="44" fill="none" stroke="{color}" stroke-width="2" opacity="0.8"/>
    <circle cx="50" cy="50" r="38" fill="none" stroke="{color}" stroke-width="0.5" opacity="0.2"/>
    <text x="50" y="44" text-anchor="middle" fill="{color}" font-size="22" font-weight="700">{i+1}</text>
    <text x="50" y="62" text-anchor="middle" fill="{color}" font-size="11" opacity="0.7">STEP</text>
  </svg>
  <div style="font-size:22px; color:#E8E5E0; margin-top:12px; font-weight:600;">{label}</div>
  <div style="font-size:14px; color:#6A7788; margin-top:4px;">{desc}</div>
</div>"""
        if i < 4:
            node_html += f"""
<div style="position:absolute; top:348px; left:{210 + i * 340}px; width:230px; height:1px; background:linear-gradient(90deg,{nodes[i][0]},{nodes[i+1][0]}); z-index:4; opacity:0.5;">
  <div style="position:absolute; right:-4px; top:-3px; width:0; height:0; border-left:7px solid {nodes[i+1][0]}; border-top:4px solid transparent; border-bottom:4px solid transparent; opacity:0.5;"></div>
</div>"""

    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.title {{ position:absolute; top:50px; left:50%; transform:translateX(-50%); font-size:36px; font-weight:700; color:#E8E5E0; letter-spacing:4px; z-index:6; }}
.title::after {{ content:''; display:block; width:50px; height:2px; background:#D4A843; margin:8px auto 0; }}
.hand-bg {{ position:absolute; left:50%; top:55%; transform:translate(-50%,-50%); z-index:1; opacity:0.06; }}
.timeline {{ position:absolute; bottom:80px; left:50%; transform:translateX(-50%); font-size:15px; color:#5A6080; z-index:5; letter-spacing:3px; }}
</style></head><body>
<div class="slide">
<div class="title">本地 CV 识别流水线 ⭐</div>
<div class="hand-bg">
  <svg width="500" height="500" viewBox="0 0 300 300">
    <path d="M110 180 Q100 140 105 100 Q108 75 118 65 Q125 60 128 68 L128 130 Q130 96 133 58 Q135 38 142 34 Q148 32 149 46 L148 118 Q152 80 155 48 Q157 28 164 26 Q170 27 168 58 L165 125 Q168 94 174 60 Q177 42 182 43 Q186 50 183 85 L176 148 Q168 188 156 210 Q144 226 130 226 Q118 225 114 210 Q108 195 110 180Z" fill="#4A7FB5" opacity="0.5"/>
  </svg>
</div>
{node_html}
<div class="timeline">输入帧 → 预处理 → 特征提取 → 分类 → 结果输出</div>
</div></body></html>"""

# ============================================================
# SLIDE 05: Fusion Strategy (KEY)
# ============================================================
def slide_05():
    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.title {{ position:absolute; top:50px; left:50%; transform:translateX(-50%); font-size:36px; font-weight:700; color:#E8E5E0; letter-spacing:4px; z-index:6; }}
.title::after {{ content:''; display:block; width:50px; height:2px; background:#D4A843; margin:8px auto 0; }}
.src {{ position:absolute; top:140px; width:260px; padding:16px 20px; border-radius:5px; text-align:center; z-index:5; }}
.src-l {{ left:320px; border:1px solid #4A7FB5; }} .src-r {{ right:320px; border:1px solid #7B6BA0; }}
.src-title {{ font-size:24px; font-weight:600; }} .st-l {{ color:#4A7FB5; }} .st-r {{ color:#7B6BA0; }}
.src-sub {{ font-size:14px; color:#7A8098; margin-top:4px; }}
.hex {{ position:absolute; top:380px; left:50%; transform:translate(-50%,-50%); z-index:5; }}
.hex svg {{ display:block; }}
.paths {{ position:absolute; top:480px; left:50%; transform:translateX(-50%); width:1100px; z-index:4; }}
.path-item {{ display:inline-block; width:340px; text-align:center; }}
.path-label {{ font-size:18px; font-weight:600; padding:8px 18px; border-radius:4px; display:inline-block; }}
.pl-green {{ color:#5BAF8A; border:1px solid #5BAF8A; }} .pl-gold {{ color:#D4A843; border:1px solid #D4A843; }} .pl-rose {{ color:#C57B8A; border:1px solid #C57B8A; }}
.path-desc {{ font-size:15px; color:#C0C4D0; margin-top:4px; }}
.path-anno {{ font-size:12px; color:#6A7088; margin-top:2px; }}
.output {{ position:absolute; bottom:70px; left:50%; transform:translateX(-50%); padding:14px 40px; border:1px solid #5BAF8A; border-radius:5px; z-index:5; font-size:18px; color:#E0E0E0; letter-spacing:2px; }}
.vl {{ position:absolute; width:1px; background:#4A5568; z-index:3; opacity:0.4; }}
</style></head><body>
<div class="slide">
<div class="title">双引擎融合策略 ⭐</div>

<div class="src src-l">
  <div class="src-title st-l">本地 CV</div>
  <div class="src-sub">~30ms / 帧 &nbsp;·&nbsp; 高频实时</div>
</div>
<div class="src src-r">
  <div class="src-title st-r">云端 AI</div>
  <div class="src-sub">按需触发 &nbsp;·&nbsp; 精准校验</div>
</div>

<div style="position:absolute; top:205px; left:450px; width:20px; height:60px; border-left:1px solid rgba(74,127,181,0.4); border-bottom:1px solid rgba(74,127,181,0.4); z-index:3;"></div>
<div style="position:absolute; top:205px; right:470px; width:20px; height:60px; border-right:1px solid rgba(123,107,160,0.4); border-bottom:1px solid rgba(123,107,160,0.4); z-index:3;"></div>

<div style="position:absolute; top:270px; left:50%; transform:translateX(-50%); width:1px; height:44px; background:rgba(74,127,181,0.4); z-index:3;"></div>
<div style="position:absolute; top:270px; left:calc(50% + 120px); width:1px; height:44px; background:rgba(123,107,160,0.4); z-index:3;"></div>

<div class="hex" style="position:absolute; top:380px; left:50%; transform:translate(-50%,-50%);">
  <svg width="200" height="120" viewBox="0 0 200 120">
    <polygon points="100,4 192,34 192,86 100,116 8,86 8,34" fill="none" stroke="#D4A843" stroke-width="2"/>
    <text x="100" y="56" text-anchor="middle" fill="#D4A843" font-size="22" font-weight="700">融合决策</text>
    <text x="100" y="76" text-anchor="middle" fill="rgba(212,168,67,0.5)" font-size="12">Fusion Engine</text>
  </svg>
</div>

<div style="position:absolute; top:440px; left:50%; transform:translateX(-50%); width:1px; height:30px; background:rgba(212,168,67,0.4); z-index:3;"></div>

<div style="position:absolute; top:480px; left:50%; transform:translateX(-50%); display:flex; gap:60px; z-index:5;">
  <div style="width:280px; text-align:center;">
    <div style="display:inline-block; padding:6px 18px; border:1px solid #5BAF8A; border-radius:4px; font-size:17px; color:#5BAF8A; font-weight:600;">OPEN_PALM</div>
    <div style="font-size:15px; color:#C0C4D0; margin-top:6px;">→ 直接采用</div>
    <div style="font-size:12px; color:#6A7088;">快速路径</div>
  </div>
  <div style="width:280px; text-align:center;">
    <div style="display:inline-block; padding:6px 18px; border:1px solid #D4A843; border-radius:4px; font-size:17px; color:#D4A843; font-weight:600;">FIST</div>
    <div style="font-size:15px; color:#C0C4D0; margin-top:6px;">→ 云端确认</div>
    <div style="font-size:12px; color:#6A7088;">安全路径</div>
  </div>
  <div style="width:280px; text-align:center;">
    <div style="display:inline-block; padding:6px 18px; border:1px solid #C57B8A; border-radius:4px; font-size:17px; color:#C57B8A; font-weight:600;">冲突</div>
    <div style="font-size:15px; color:#C0C4D0; margin-top:6px;">→ 放弃此帧</div>
    <div style="font-size:12px; color:#6A7088;">兜底策略</div>
  </div>
</div>

<div style="position:absolute; bottom:85px; left:50%; transform:translateX(-50%); display:flex; gap:60px; z-index:4;">
  <div style="width:1px; height:24px; background:rgba(91,175,138,0.3);"></div>
  <div style="width:1px; height:24px; background:rgba(212,168,67,0.3);"></div>
  <div style="width:1px; height:24px; background:rgba(197,123,138,0.3);"></div>
</div>

<div class="output">
  GestureCommand &nbsp;→&nbsp; 发送至游戏层
</div>
</div></body></html>"""

# ============================================================
# SLIDE 06: Design Patterns (KEY)
# ============================================================
def slide_06():
    cards = [
        ("#4A7FB5", "Strategy", "策略模式", "识别器可动态切换", "本地 CV ↔ 云端 AI"),
        ("#D4A843", "Singleton", "单例模式", "全局唯一融合器", "线程安全 · 延迟初始化"),
        ("#5BAF8A", "Observer", "观察者模式", "帧推送与回调解耦", "非阻塞事件通知"),
    ]
    cards_html = ""
    for i, (color, eng, cn, desc, anno) in enumerate(cards):
        cards_html += f"""
<div style="width:300px; background:rgba(30,35,64,0.9); border-radius:6px; padding:30px 22px 24px; box-shadow:0 6px 22px rgba(0,0,0,0.4); text-align:center;">
  <div style="height:3px; width:55%; background:{color}; margin:0 auto 20px; border-radius:2px;"></div>
  <div style="font-size:34px; font-weight:700; color:{color}; margin-bottom:6px;">{eng}</div>
  <div style="font-size:20px; color:#D8D5D0; margin-bottom:16px;">{cn}</div>
  <div style="width:40px; height:1px; background:rgba(255,255,255,0.1); margin:0 auto 16px;"></div>
  <div style="font-size:16px; color:#D4A843; margin-bottom:8px;">{desc}</div>
  <div style="font-size:12px; color:#6A7088;">{anno}</div>
</div>"""

    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.title {{ position:absolute; top:50px; left:50%; transform:translateX(-50%); font-size:36px; font-weight:700; color:#E8E5E0; letter-spacing:4px; z-index:6; }}
.title::after {{ content:''; display:block; width:50px; height:2px; background:#D4A843; margin:8px auto 0; }}
.footnote {{ position:absolute; bottom:70px; left:50%; transform:translateX(-50%); font-size:14px; color:#6A7088; z-index:5; letter-spacing:2px; }}
</style></head><body>
<div class="slide">
<div class="title">设计模式应用 ⭐</div>
<div style="position:absolute; top:160px; left:50%; transform:translateX(-50%); display:flex; gap:40px; z-index:5;">
{cards_html}
</div>
<div class="footnote">三种模式协同：策略切换 → 单例融合 → 观察者推送</div>
</div></body></html>"""

# ============================================================
# SLIDE 07: Architecture
# ============================================================
def slide_07():
    layers = [
        ("#D4A843", "游戏层", ["Game", "Shop", "Menu"], "Cocos2d-x 引擎"),
        ("#4A7FB5", "融合层", ["GestureFusion"], "决策中间件"),
        ("#5BAF8A", "识别层", ["Local CV", "Cloud AI"], "双识别引擎"),
        ("#7B6BA0", "云端服务", ["AI API", "HTTP"], "外部服务"),
    ]
    layers_html = ""
    for i, (color, name, tags, _) in enumerate(layers):
        top = 160 + i * 130
        tags_html = "".join([f'<span style="padding:4px 14px; border:1px solid {color}; border-radius:3px; font-size:14px; color:{color}; margin-left:10px;">{t}</span>' for t in tags])
        layers_html += f"""
<div style="position:absolute; top:{top}px; left:120px; right:120px; height:100px; background:rgba({','.join(str(int(color[i:i+2],16)) for i in (1,3,5))},0.06); border-radius:4px; z-index:4; display:flex; align-items:center; justify-content:space-between; padding:0 40px;">
  <div style="font-size:28px; font-weight:600; color:#E0DDD8;">{name}</div>
  <div>{tags_html}</div>
</div>"""
        if i < 3:
            layers_html += f"""
<div style="position:absolute; top:{top+100}px; left:180px; z-index:3; display:flex; align-items:center; gap:8px;">
  <div style="width:1px; height:28px; background:rgba({','.join(str(int(color[i:i+2],16)) for i in (1,3,5))},0.3);"></div>
  <span style="font-size:12px; color:#6A7088;">↓ {["帧数据","识别请求","API 调用"][i]}</span>
</div>"""

    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.title {{ position:absolute; top:50px; left:50%; transform:translateX(-50%); font-size:38px; font-weight:700; color:#E8E5E0; letter-spacing:4px; z-index:6; }}
.title::after {{ content:''; display:block; width:50px; height:2px; background:#D4A843; margin:8px auto 0; }}
.footnote {{ position:absolute; bottom:65px; left:50%; transform:translateX(-50%); font-size:13px; color:#5A6080; z-index:5; }}
</style></head><body>
<div class="slide">
<div class="title">系统架构</div>
{layers_html}
<div class="footnote">Cocos2d-x 游戏引擎 ← OpenCV 本地识别 ← 云端 AI 服务</div>
</div></body></html>"""

# ============================================================
# SLIDE 08: Testing
# ============================================================
def slide_08():
    metrics = [
        ("35", "FPS", "≥ 30"),
        ("85", "ms", "≤ 100"),
        ("92", "% 准确率", "≥ 90%"),
    ]
    cards_html = ""
    for val, unit, target in metrics:
        cards_html += f"""
<div style="width:280px; background:rgba(30,35,64,0.9); border-radius:6px; padding:32px 20px; box-shadow:0 6px 22px rgba(0,0,0,0.4); text-align:center;">
  <div style="font-size:80px; font-weight:900; color:#5BAF8A; line-height:1;">{val}</div>
  <div style="font-size:22px; color:#E0DDD8; margin-top:6px;">{unit}</div>
  <div style="width:40px; height:1px; background:rgba(255,255,255,0.08); margin:16px auto;"></div>
  <div style="font-size:15px; color:#7A8098;">目标 {target}</div>
  <div style="font-size:22px; color:#5BAF8A; margin-top:6px;">✓</div>
</div>"""

    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.title {{ position:absolute; top:50px; left:50%; transform:translateX(-50%); font-size:38px; font-weight:700; color:#E8E5E0; letter-spacing:4px; z-index:6; }}
.title::after {{ content:''; display:block; width:50px; height:2px; background:#D4A843; margin:8px auto 0; }}
.bar {{ position:absolute; bottom:90px; left:50%; transform:translateX(-50%); width:800px; height:3px; background:rgba(255,255,255,0.06); border-radius:2px; z-index:5; }}
.bar-fill {{ height:100%; background:#5BAF8A; border-radius:2px; }}
.bar-label {{ position:absolute; bottom:63px; left:50%; transform:translateX(-50%); font-size:13px; color:#6A7088; z-index:5; }}
</style></head><body>
<div class="slide">
<div class="title">测试与性能</div>
<div style="position:absolute; top:200px; left:50%; transform:translateX(-50%); display:flex; gap:50px; z-index:5;">
{cards_html}
</div>
<div class="bar"><div class="bar-fill" style="width:100%;"></div></div>
<div class="bar-label">全部指标达标 ✓</div>
</div></body></html>"""

# ============================================================
# SLIDE 09: Innovation (KEY)
# ============================================================
def slide_09():
    items = [
        ("#4A7FB5", "双引擎融合", "本地 CV 快速响应 + 云端 AI 精准校验", "快速路径与安全路径智能分流"),
        ("#5BAF8A", "手势即指令", "无需键盘鼠标，直觉式自然交互", "OPEN_PALM 瞄准 / FIST 释放"),
        ("#4A7FB5", "异步非阻塞", "帧推送与识别回调完全解耦", "单帧处理延迟 ≤ 100ms"),
        ("#C57B8A", "状态机去抖", "时域滤波防止手势误触发", "连续 N 帧一致才确认触发"),
    ]
    cards = ""
    for i, (color, title, desc, anno) in enumerate(items):
        top = 180 if i < 2 else 480
        left = 120 if i % 2 == 0 else 620
        cards += f"""
<div style="position:absolute; top:{top}px; left:{left}px; width:540px; background:rgba(30,35,64,0.9); border-radius:6px; padding:32px 30px; box-shadow:0 6px 22px rgba(0,0,0,0.4);">
  <div style="height:3px; width:50px; background:{color}; margin-bottom:18px; border-radius:2px;"></div>
  <div style="font-size:30px; font-weight:700; color:#E0DDD8; margin-bottom:8px;">{title}</div>
  <div style="font-size:17px; color:#D4A843; margin-bottom:6px;">{desc}</div>
  <div style="font-size:13px; color:#6A7088;">{anno}</div>
</div>"""

    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.title {{ position:absolute; top:50px; left:50%; transform:translateX(-50%); font-size:40px; font-weight:700; color:#D4A843; letter-spacing:6px; z-index:6; }}
.title::after {{ content:''; display:block; width:60px; height:2px; background:#D4A843; margin:8px auto 0; }}
</style></head><body>
<div class="slide">
<div class="title">创新亮点 ⭐</div>
{cards}
<div style="position:absolute; top:425px; left:50%; transform:translateX(-50%); width:6px; height:6px; background:#D4A843; border-radius:50%; z-index:6; opacity:0.5;"></div>
</div></body></html>"""

# ============================================================
# SLIDE 10: Closing
# ============================================================
def slide_10():
    return f"""<!DOCTYPE html><html lang="zh"><head><meta charset="UTF-8"><style>
{CSS_COMMON}
.hook {{ position:absolute; top:80px; left:50%; transform:translateX(-50%); z-index:5; }}
.thanks {{ font-family:'Georgia','SimSun',serif; font-size:140px; font-weight:900;
  background:linear-gradient(180deg,#F0D68A 0%,#D4A843 35%,#B8860B 60%,#E8C870 80%,#D4A843 100%);
  -webkit-background-clip:text; -webkit-text-fill-color:transparent; background-clip:text;
  filter:drop-shadow(0 4px 10px rgba(0,0,0,0.5)); text-align:center; }}
.thanks-area {{ position:absolute; top:380px; left:50%; transform:translateX(-50%); text-align:center; z-index:6; }}
.sub-text {{ font-size:24px; color:#C0C4D0; margin-top:20px; letter-spacing:4px; }}
.github {{ font-size:18px; color:#4A7FB5; margin-top:18px; font-family:'Consolas',monospace; letter-spacing:1px; }}
.author {{ font-size:16px; color:#6A7088; margin-top:40px; }}
.nuggets {{ position:absolute; bottom:50px; left:50%; transform:translateX(-50%); display:flex; gap:12px; align-items:flex-end; z-index:5; }}
.nuggets img {{ display:block; filter:drop-shadow(0 3px 10px rgba(0,0,0,0.5)); }}
</style></head><body>
<div class="slide">
<div class="hook">
  <svg width="120" height="130" viewBox="0 0 120 130">
    <defs><linearGradient id="hg2" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0%" stop-color="#D4A843"/><stop offset="40%" stop-color="#F0D88A"/>
      <stop offset="65%" stop-color="#B8860B"/><stop offset="100%" stop-color="#8B6914"/>
    </linearGradient></defs>
    <circle cx="60" cy="14" r="8" fill="none" stroke="url(#hg2)" stroke-width="3"/>
    <path d="M60 22 L60 55 Q60 82 42 96 Q30 104 35 120 Q40 130 55 124 Q64 120 61 108 Q58 95 45 88" fill="none" stroke="url(#hg2)" stroke-width="5" stroke-linecap="round"/>
    <path d="M35 120 L28 133" fill="none" stroke="url(#hg2)" stroke-width="4" stroke-linecap="round"/>
  </svg>
</div>
<div class="thanks-area">
  <div class="thanks">谢谢</div>
  <div class="sub-text">黄金矿工 — 手势识别版</div>
  <div class="github">github.com/mzk-C4/gold-miner-vision</div>
  <div class="author">mzk-C4</div>
</div>
<div class="nuggets">
  <img src="{SPRITES['gold-nugget-2']}" style="height:55px">
  <img src="{SPRITES['diamond-2']}" style="height:50px">
  <img src="{SPRITES['gold-nugget-1']}" style="height:42px">
  <img src="{SPRITES['gen-diamond']}" style="height:52px">
  <img src="{SPRITES['gold-nugget-3']}" style="height:48px">
</div>
</div></body></html>"""

# ============================================================
# MAIN: Generate all 10 slides
# ============================================================
SLIDES = [
    ("01-封面", slide_01),
    ("02-项目背景", slide_02),
    ("03-核心数据结构", slide_03),
    ("04-CV识别流水线", slide_04),
    ("05-双引擎融合策略", slide_05),
    ("06-设计模式应用", slide_06),
    ("07-系统架构", slide_07),
    ("08-测试与性能", slide_08),
    ("09-创新亮点", slide_09),
    ("10-结尾", slide_10),
]

def main():
    print("=" * 60)
    print("Generating 10-slide PPTX presentation")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = PPTX_W
    prs.slide_height = PPTX_H
    blank_layout = prs.slide_layouts[6]

    with sync_playwright() as p:
        browser = p.chromium.launch()

        for slide_num, (name, slide_func) in enumerate(SLIDES, 1):
            print(f"\n[{slide_num}/10] {name}")

            print("  Building HTML...")
            html = slide_func()

            print("  Screenshotting...")
            screenshot = render_slide(html, browser)

            print("  Adding to PPTX...")
            slide = prs.slides.add_slide(blank_layout)
            img_stream = io.BytesIO(screenshot)
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0), PPTX_W, PPTX_H)

            print(f"  [OK] Done")

        browser.close()

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))

    file_size = OUTPUT_PPTX.stat().st_size / (1024 * 1024)
    print(f"\n{'=' * 60}")
    print(f"Complete! 10 slides saved to:")
    print(f"  {OUTPUT_PPTX}")
    print(f"  File size: {file_size:.1f} MB")
    print(f"{'=' * 60}")

if __name__ == "__main__":
    main()
