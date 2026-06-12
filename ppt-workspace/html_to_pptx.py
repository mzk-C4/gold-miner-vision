"""
Convert HTML slides to PPTX via Playwright screenshots.
Usage: python html_to_pptx.py
Output: slide-01-cover.pptx (and later, combined PPTX with all slides)
"""

import os
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

WORKSPACE = Path(r"D:\GOLD MINER\ppt-workspace")
HTML_DIR = WORKSPACE / "html"
PPTX_DIR = WORKSPACE / "pptx"
PPTX_DIR.mkdir(exist_ok=True)

# Slide dimensions (16:9 at 1920x1080 → 13.333" x 7.5")
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def html_to_image(html_path: Path, output_path: Path):
    """Screenshot HTML file at 1920x1080 using Playwright."""
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.goto(f"file:///{html_path.as_posix()}")
        # Wait for fonts and SVG to render
        page.wait_for_timeout(1500)
        page.screenshot(path=str(output_path), full_page=False)
        browser.close()
    print(f"  Screenshot saved: {output_path.name}")
    return output_path

def image_to_pptx_slide(image_path: Path, output_pptx: Path):
    """Create a PPTX with the image as a full-slide background."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add image as full-slide background
    slide.shapes.add_picture(
        str(image_path),
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )

    prs.save(str(output_pptx))
    print(f"  PPTX saved: {output_pptx.name}")

def process_slide(html_name: str):
    """Process a single HTML → PNG → PPTX pipeline."""
    html_path = HTML_DIR / html_name
    if not html_path.exists():
        print(f"  SKIP: {html_name} not found")
        return None

    base_name = html_name.replace(".html", "")
    png_path = PPTX_DIR / f"{base_name}.png"
    pptx_path = PPTX_DIR / f"{base_name}.pptx"

    print(f"\nProcessing: {html_name}")
    html_to_image(html_path, png_path)
    image_to_pptx_slide(png_path, pptx_path)
    return pptx_path

def combine_all_to_pptx(html_files: list, output_name: str):
    """Combine multiple HTML slides into a single PPTX."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    with sync_playwright() as p:
        browser = p.chromium.launch()
        for html_file in html_files:
            html_path = HTML_DIR / html_file
            if not html_path.exists():
                print(f"  SKIP: {html_file} not found")
                continue

            print(f"  Screenshotting: {html_file}")
            page = browser.new_page(viewport={"width": 1920, "height": 1080})
            page.goto(f"file:///{html_path.as_posix()}")
            page.wait_for_timeout(1500)

            # Screenshot to memory
            screenshot_bytes = page.screenshot(full_page=False)
            page.close()

            # Add slide with image
            slide = prs.slides.add_slide(blank_layout)
            image_stream = io.BytesIO(screenshot_bytes)
            slide.shapes.add_picture(
                image_stream,
                Inches(0), Inches(0),
                SLIDE_WIDTH, SLIDE_HEIGHT
            )
            print(f"    Added to PPTX")

        browser.close()

    output_path = PPTX_DIR / output_name
    prs.save(str(output_path))
    print(f"\nCombined PPTX saved: {output_path}")
    return output_path


if __name__ == "__main__":
    # Find all HTML slides
    html_files = sorted([f.name for f in HTML_DIR.glob("slide-*.html")])

    if not html_files:
        print("No HTML slides found in", HTML_DIR)
        exit(1)

    print(f"Found {len(html_files)} HTML slide(s)")
    print("=" * 50)

    # Process individually
    for html_file in html_files:
        process_slide(html_file)

    # Combine into one PPTX if multiple slides
    if len(html_files) > 1:
        print("\n" + "=" * 50)
        print("Creating combined PPTX...")
        combine_all_to_pptx(html_files, "黄金矿工-手势识别版.pptx")
    elif len(html_files) == 1:
        print("\nSingle slide — individual PPTX ready.")

    print("\nDone! All files in:", PPTX_DIR)
